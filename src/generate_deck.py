import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    bg_dark = RGBColor(6, 11, 19)        # #060b13
    primary_color = RGBColor(16, 185, 129)  # #10b981 (Mint Green)
    secondary_color = RGBColor(59, 130, 246) # #3b82f6 (Electric Blue)
    text_light = RGBColor(241, 245, 249)     # #f1f5f9
    text_muted = RGBColor(148, 163, 184)     # #94a3b8
    
    # Helper to set slide background
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_dark
        
    # Helper to add a slide title
    def add_slide_header(slide, title_text, category_text="EV TARIFF OPTIMIZATION"):
        # Category/Tracker
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.name = 'Outfit'
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = primary_color
        
        # Main Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.name = 'Outfit'
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = text_light

    # ==========================================
    # SLIDE 1: Title Slide (Cover)
    # ==========================================
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide1 = prs.slides.add_slide(slide_layout)
    set_background(slide1)
    # Main Title & Subtitle Box
    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_top = Inches(0)
    
    p_title = tf.paragraphs[0]
    p_title.text = "Agentic AI-Based Dynamic Tariff Optimization"
    p_title.font.name = 'Outfit'
    p_title.font.size = Pt(42)
    p_title.font.bold = True
    p_title.font.color.rgb = primary_color
    p_title.space_after = Pt(12)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "For EV Charging Networks Using Large-Scale Charging Session Data"
    p_sub.font.name = 'Outfit'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = text_light
    
    # Submitter & Organization Box (aligned below title box)
    authorBox = slide1.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(11.3), Inches(1.8))
    tf_author = authorBox.text_frame
    tf_author.word_wrap = True
    tf_author.margin_left = Inches(0)
    tf_author.margin_top = Inches(0)
    
    p_org = tf_author.paragraphs[0]
    p_org.text = "Open Project 2026 Submission by Society of Business"
    p_org.font.name = 'Inter'
    p_org.font.size = Pt(14)
    p_org.font.italic = True
    p_org.font.color.rgb = text_muted
    p_org.space_after = Pt(8)
    
    p_author = tf_author.add_paragraph()
    p_author.text = "Submitted by: Aryan Mehra (Enrollment: 23115025)"
    p_author.font.name = 'Inter'
    p_author.font.size = Pt(15)
    p_author.font.bold = True
    p_author.font.color.rgb = primary_color

    # ==========================================
    # SLIDE 2: Project Overview & Objectives
    # ==========================================
    slide2 = prs.slides.add_slide(slide_layout)
    set_background(slide2)
    add_slide_header(slide2, "Project Overview & Objectives", "PROJECT OVERVIEW")
    
    # Left Column / Bullet Box
    leftBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.8))
    tf_left2 = leftBox2.text_frame
    tf_left2.word_wrap = True
    
    p_lh2 = tf_left2.paragraphs[0]
    p_lh2.text = "Dynamic Tariff Optimization Problem Setup"
    p_lh2.font.name = 'Outfit'
    p_lh2.font.size = Pt(20)
    p_lh2.font.bold = True
    p_lh2.font.color.rgb = secondary_color
    p_lh2.space_after = Pt(14)
    
    bullets_overview = [
        "Problem Statement: Flat-rate tariffs cause peak-hour congestion and underutilization in off-peak windows.",
        "Objective: Build a self-improving pricing engine that autonomously predicts demand, optimizes tariffs, and learns from outcomes to maximize revenue and balance grid load.",
        "Datasets: Caltech ACN-Data (30,000+ sessions) and Shenzhen UrbanEV Dataset (24,798 charging piles)."
    ]
    for bullet in bullets_overview:
        p = tf_left2.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 3: Preprocessing
    # ==========================================
    slide3 = prs.slides.add_slide(slide_layout)
    set_background(slide3)
    add_slide_header(slide3, "Data Preprocessing & Feature Engineering")
    
    # Left Column: Preprocessing Decisions
    leftBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left3 = leftBox3.text_frame
    tf_left3.word_wrap = True
    
    p_lh3 = tf_left3.paragraphs[0]
    p_lh3.text = "Data Cleaning & Alignment"
    p_lh3.font.name = 'Outfit'
    p_lh3.font.size = Pt(20)
    p_lh3.font.bold = True
    p_lh3.font.color.rgb = secondary_color
    p_lh3.space_after = Pt(14)
    
    bullets_left = [
        "ACN Dataset: Cleaned session data and forward-filled charging metrics into nested user update sheets.",
        "ACN Dataset: Aggregated session-level metrics into hourly active occupancies and charging loads.",
        "UrbanEV: Reshaped 2.13M row grid matrix datasets into long format for temporal-spatial training.",
        "Unified Grid Scaling: Multiplied default price by 15.0 to align prices with Indian Rupee baseline metrics (average ~₹14.38/kWh)."
    ]
    for bullet in bullets_left:
        p = tf_left3.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Feature Engineering
    rightBox3 = slide3.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_right3 = rightBox3.text_frame
    tf_right3.word_wrap = True
    
    p_rh3 = tf_right3.paragraphs[0]
    p_rh3.text = "Feature Engineering Decisions"
    p_rh3.font.name = 'Outfit'
    p_rh3.font.size = Pt(20)
    p_rh3.font.bold = True
    p_rh3.font.color.rgb = secondary_color
    p_rh3.space_after = Pt(14)
    
    bullets_right = [
        "Charger Utilization Rate: Computed as active charging time divided by total available capacity.",
        "Queue Length Proxy: Defined as max(0, Occupancy - Capacity) to estimate queue congestion at peak intervals.",
        "Occupancy Density: Normalized grid occupancies per square kilometer.",
        "Time-Series Lags: Engineered 1-hour, 2-hour, and 24-hour lags, alongside 3-hour rolling averages, to capture temporal patterns."
    ]
    for bullet in bullets_right:
        p = tf_right3.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)

    # ==========================================
    # SLIDE 4: Exploratory Data Analysis (EDA)
    # ==========================================
    slide4 = prs.slides.add_slide(slide_layout)
    set_background(slide4)
    add_slide_header(slide4, "Exploratory Data Analysis (EDA)")
    
    leftBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left4 = leftBox4.text_frame
    tf_left4.word_wrap = True
    
    p_lh4 = tf_left4.paragraphs[0]
    p_lh4.text = "Temporal & Spatial Insights"
    p_lh4.font.name = 'Outfit'
    p_lh4.font.size = Pt(20)
    p_lh4.font.bold = True
    p_lh4.font.color.rgb = secondary_color
    p_lh4.space_after = Pt(14)
    
    bullets_eda = [
        "Workplace Patterns (ACN Caltech): Strong charging demand peaks during weekday business hours (8 AM - 4 PM), collapsing to near-zero on weekends.",
        "Public Urban Grids (UrbanEV Shenzhen): Bimodal demand peaks around morning commute (8 AM - 10 AM) and early evening (6 PM - 9 PM) as taxis and public fleets recharge.",
        "Congestion Inefficiency: Flat baseline pricing (fixed ₹15/kWh) fails to represent operational strain, leading to severe localized queues during peak hours while chargers sit idle off-peak."
    ]
    for bullet in bullets_eda:
        p = tf_left4.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Visualizations & Analysis Note / Image fallback
    import os
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    eda_img_path = os.path.join(project_dir, "outputs", "eda_profile.png")
    if os.path.exists(eda_img_path):
        slide4.shapes.add_picture(eda_img_path, Inches(6.8), Inches(1.8), Inches(5.6), Inches(3.7))
    else:
        rightBox4 = slide4.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
        tf_right4 = rightBox4.text_frame
        tf_right4.word_wrap = True
        
        p_rh4 = tf_right4.paragraphs[0]
        p_rh4.text = "Pricing Integration Hook"
        p_rh4.font.name = 'Outfit'
        p_rh4.font.size = Pt(20)
        p_rh4.font.bold = True
        p_rh4.font.color.rgb = secondary_color
        p_rh4.space_after = Pt(14)
        
        p_desc4 = tf_right4.add_paragraph()
        p_desc4.text = "An interactive visual chart showing average daily profiles is embedded in the web dashboard console, highlighting the peak volatility between weekday work commutes and evening grid peaks."
        p_desc4.font.name = 'Inter'
        p_desc4.font.size = Pt(14)
        p_desc4.font.color.rgb = text_light
        p_desc4.space_after = Pt(14)
        
        p_desc4_sub = tf_right4.add_paragraph()
        p_desc4_sub.text = "Grid demand is highly elastic if priced correctly. Shifts in pricing directly affect charger utilization rate during off-peak windows."
        p_desc4_sub.font.name = 'Inter'
        p_desc4_sub.font.size = Pt(13)
        p_desc4_sub.font.italic = True
        p_desc4_sub.font.color.rgb = text_muted

    # ==========================================
    # SLIDE 5: Demand Forecasting Agent
    # ==========================================
    slide5 = prs.slides.add_slide(slide_layout)
    set_background(slide5)
    add_slide_header(slide5, "Demand Prediction Agent")
    
    leftBox5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left5 = leftBox5.text_frame
    tf_left5.word_wrap = True
    
    p_lh5 = tf_left5.paragraphs[0]
    p_lh5.text = "Forecasting Methodology"
    p_lh5.font.name = 'Outfit'
    p_lh5.font.size = Pt(20)
    p_lh5.font.bold = True
    p_lh5.font.color.rgb = secondary_color
    p_lh5.space_after = Pt(14)
    
    bullets_pred = [
        "LightGBM Regressors: Trained separate models to predict charging occupancy and volume (expected load) simultaneously.",
        "Lag & Rolling Windows: Utilized temporal features (hour, day, weekend markers) alongside historical load lags to capture auto-regressive properties.",
        "Congestion Probability: Calculated using a sigmoid activation function centered around 80% utilization rate, warning operators of overload risks in real time."
    ]
    for bullet in bullets_pred:
        p = tf_left5.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Results & Model Performance Metrics
    rightBox5 = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_right5 = rightBox5.text_frame
    tf_right5.word_wrap = True
    
    p_rh5 = tf_right5.paragraphs[0]
    p_rh5.text = "Model Evaluation Performance"
    p_rh5.font.name = 'Outfit'
    p_rh5.font.size = Pt(20)
    p_rh5.font.bold = True
    p_rh5.font.color.rgb = secondary_color
    p_rh5.space_after = Pt(14)
    
    metrics = [
        ("UrbanEV Occupancy R²", "98.68%", "Highly predictive spatial-temporal consistency"),
        ("UrbanEV Charging Volume R²", "91.07%", "Robust regression over aggregated loads"),
        ("ACN Occupancy R²", "95.24%", "Predictive accuracy over workplace schedules"),
        ("ACN Charging Volume R²", "80.57%", "Captures high-volatility session curves")
    ]
    for metric_name, score, detail in metrics:
        p_name = tf_right5.add_paragraph()
        p_name.text = f"{metric_name}: "
        p_name.font.name = 'Inter'
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = text_light
        
        # Highlight score
        run = p_name.add_run()
        run.text = score
        run.font.bold = True
        run.font.color.rgb = primary_color
        
        p_det = tf_right5.add_paragraph()
        p_det.text = f"  - {detail}"
        p_det.font.name = 'Inter'
        p_det.font.size = Pt(11)
        p_det.font.color.rgb = text_muted
        p_det.space_after = Pt(8)

    # ==========================================
    # SLIDE 6: Tariff Pricing Agent
    # ==========================================
    slide6 = prs.slides.add_slide(slide_layout)
    set_background(slide6)
    add_slide_header(slide6, "Tariff Pricing Agent")
    
    leftBox6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left6 = leftBox6.text_frame
    tf_left6.word_wrap = True
    
    p_lh6 = tf_left6.paragraphs[0]
    p_lh6.text = "Dynamic Tariff Logic"
    p_lh6.font.name = 'Outfit'
    p_lh6.font.size = Pt(20)
    p_lh6.font.bold = True
    p_lh6.font.color.rgb = secondary_color
    p_lh6.space_after = Pt(14)
    
    bullets_tariff = [
        "Surge Pricing Boundary: Triggered when predicted station utilization exceeds 80%. Linearly scales tariff from baseline (₹15/kWh) up to max cap (₹25/kWh).",
        "Discount Pricing Boundary: Triggered when predicted station utilization drops below 30%. Linearly decreases tariff from baseline down to floor (₹10/kWh).",
        "Fixed Baseline Pricing: Holds pricing flat at ₹15/kWh for normal occupancy states (30% to 80% utilization)."
    ]
    for bullet in bullets_tariff:
        p = tf_left6.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Elasticity Demand Response Simulation / Image fallback
    sim_img_path = os.path.join(project_dir, "outputs", "simulation_results.png")
    if os.path.exists(sim_img_path):
        rightBox6 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(1.8))
        tf_right6 = rightBox6.text_frame
        tf_right6.word_wrap = True
        tf_right6.margin_left = Inches(0)
        tf_right6.margin_top = Inches(0)
        
        p_rh6 = tf_right6.paragraphs[0]
        p_rh6.text = "Simulating User Demand Elasticity"
        p_rh6.font.name = 'Outfit'
        p_rh6.font.size = Pt(20)
        p_rh6.font.bold = True
        p_rh6.font.color.rgb = secondary_color
        p_rh6.space_after = Pt(8)
        
        p_eq = tf_right6.add_paragraph()
        p_eq.text = "D_elastic = D_predicted * (1 - elasticity * (P_dynamic - P_baseline) / P_baseline)"
        p_eq.font.name = 'Courier New'
        p_eq.font.size = Pt(12)
        p_eq.font.bold = True
        p_eq.font.color.rgb = primary_color
        
        slide6.shapes.add_picture(sim_img_path, Inches(6.8), Inches(3.2), Inches(5.6), Inches(3.5))
    else:
        rightBox6 = slide6.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
        tf_right6 = rightBox6.text_frame
        tf_right6.word_wrap = True
        
        p_rh6 = tf_right6.paragraphs[0]
        p_rh6.text = "Simulating User Demand Elasticity"
        p_rh6.font.name = 'Outfit'
        p_rh6.font.size = Pt(20)
        p_rh6.font.bold = True
        p_rh6.font.color.rgb = secondary_color
        p_rh6.space_after = Pt(14)
        
        p_formula = tf_right6.add_paragraph()
        p_formula.text = "Demand response is modeled using the price elasticity of demand coefficient:"
        p_formula.font.name = 'Inter'
        p_formula.font.size = Pt(13)
        p_formula.font.color.rgb = text_light
        p_formula.space_after = Pt(8)
        
        p_eq = tf_right6.add_paragraph()
        p_eq.text = "D_elastic = D_predicted * (1 - elasticity * (P_dynamic - P_baseline) / P_baseline)"
        p_eq.font.name = 'Courier New'
        p_eq.font.size = Pt(13)
        p_eq.font.bold = True
        p_eq.font.color.rgb = primary_color
        p_eq.space_after = Pt(14)
        
        bullets_elast = [
            "Peak Shaving: Incentivizes users with flexible schedules to delay charging during high surge cost windows.",
            "Off-peak Filling: Stimulates off-peak demand by offering dynamic discount tariffs."
        ]
        for bullet in bullets_elast:
            p = tf_right6.add_paragraph()
            p.text = "• " + bullet
            p.font.name = 'Inter'
            p.font.size = Pt(13)
            p.font.color.rgb = text_light
            p.space_after = Pt(8)

    # ==========================================
    # SLIDE 7: Monitoring & Learning Agent
    # ==========================================
    slide7 = prs.slides.add_slide(slide_layout)
    set_background(slide7)
    add_slide_header(slide7, "Monitoring & Learning Agent")
    
    leftBox7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left7 = leftBox7.text_frame
    tf_left7.word_wrap = True
    
    p_lh7 = tf_left7.paragraphs[0]
    p_lh7.text = "Closed-Loop Feedback Loop"
    p_lh7.font.name = 'Outfit'
    p_lh7.font.size = Pt(20)
    p_lh7.font.bold = True
    p_lh7.font.color.rgb = secondary_color
    p_lh7.space_after = Pt(14)
    
    bullets_monitor = [
        "Continuous Evaluation: Tracks operational metrics (Revenue Gain %, Off-Peak Uplift, Pricing Efficiency, and Queue Reduction) across simulated daily episodes.",
        "Elasticity Tuning: Automatically adjusts pricing elasticity coefficient dynamically to avoid over-corrections and optimize revenue.",
        "Risk Mitigation: Decreases elasticity bounds automatically if a severe revenue drop occurs (e.g. drop > 10% vs baseline), preventing customer churn."
    ]
    for bullet in bullets_monitor:
        p = tf_left7.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Simulation Outcomes
    rightBox7 = slide7.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_right7 = rightBox7.text_frame
    tf_right7.word_wrap = True
    
    p_rh7 = tf_right7.paragraphs[0]
    p_rh7.text = "Simulation Outcomes (Shenzhen)"
    p_rh7.font.name = 'Outfit'
    p_rh7.font.size = Pt(20)
    p_rh7.font.bold = True
    p_rh7.font.color.rgb = secondary_color
    p_rh7.space_after = Pt(14)
    
    results = [
        ("Off-Peak Uplift", "+2.09%", "Low-demand hours volume increased"),
        ("Customer Response Rate", "0.229", "Elasticity factor observed in test splits"),
        ("Avg Dynamic Price", "₹14.55/kWh", "Lower than fixed ₹15 baseline (consumer benefit)"),
        ("Revenue Gain", "-2.27%", "Balanced grid stabilization & user relief")
    ]
    for r_name, r_val, r_det in results:
        p_name = tf_right7.add_paragraph()
        p_name.text = f"{r_name}: "
        p_name.font.name = 'Inter'
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = text_light
        
        run = p_name.add_run()
        run.text = r_val
        run.font.bold = True
        run.font.color.rgb = primary_color
        
        p_det = tf_right7.add_paragraph()
        p_det.text = f"  - {r_det}"
        p_det.font.name = 'Inter'
        p_det.font.size = Pt(11)
        p_det.font.color.rgb = text_muted
        p_det.space_after = Pt(8)

    # ==========================================
    # SLIDE 8: Business & Policy Implications
    # ==========================================
    slide8 = prs.slides.add_slide(slide_layout)
    set_background(slide8)
    add_slide_header(slide8, "Business & Grid Implications")
    
    leftBox8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_left8 = leftBox8.text_frame
    tf_left8.word_wrap = True
    
    p_lh8 = tf_left8.paragraphs[0]
    p_lh8.text = "Grid & Infrastructure Benefits"
    p_lh8.font.name = 'Outfit'
    p_lh8.font.size = Pt(20)
    p_lh8.font.bold = True
    p_lh8.font.color.rgb = secondary_color
    p_lh8.space_after = Pt(14)
    
    bullets_grid = [
        "Grid Transformer Overload Mitigation: Shaving peak charging load by up to 15% directly decreases utility infrastructure strain and distribution failures.",
        "Increased Charger Asset ROI: Incentivizing off-peak usage spreads demand, maximizing the utilization rate of expensive rapid charging assets.",
        "Wait-Time & Queuing Reduction: Distributed session times mean shorter queues during peak hours, yielding better driver satisfaction."
    ]
    for bullet in bullets_grid:
        p = tf_left8.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)
        
    # Right Column: Policy & Sustainable Pricing
    rightBox8 = slide8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_right8 = rightBox8.text_frame
    tf_right8.word_wrap = True
    
    p_rh8 = tf_right8.paragraphs[0]
    p_rh8.text = "Policy & Sustainability Alignment"
    p_rh8.font.name = 'Outfit'
    p_rh8.font.size = Pt(20)
    p_rh8.font.bold = True
    p_rh8.font.color.rgb = secondary_color
    p_rh8.space_after = Pt(14)
    
    bullets_policy = [
        "Green Solar Matching: Allows station operators to map dynamic discount pricing directly to regional solar generation peak hours (10 AM - 3 PM).",
        "Encouraging Clean EV Adoption: Dynamic discounts help keep the cost per kilometer low, incentivizing fleet operators (e.g. public taxi fleets) to electrify.",
        "Operational Transparency: Ensures price-sensitivity feedback loop parameters are auditable and adaptable to changing seasonal demand curves."
    ]
    for bullet in bullets_policy:
        p = tf_right8.add_paragraph()
        p.text = "• " + bullet
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = text_light
        p.space_after = Pt(12)

    # Save presentation
    import os
    script_dir = os.path.dirname(os.path.abspath(__file__))
    project_dir = os.path.dirname(script_dir)
    
    # Save to outputs/
    outputs_dir = os.path.join(project_dir, "outputs")
    os.makedirs(outputs_dir, exist_ok=True)
    output_path_out = os.path.join(outputs_dir, "OP26_Analytics_Presentation.pptx")
    prs.save(output_path_out)
    print(f"Presentation saved successfully to {output_path_out}")
    
    # Save to app/ for browser download
    app_dir = os.path.join(project_dir, "app")
    os.makedirs(app_dir, exist_ok=True)
    output_path_app = os.path.join(app_dir, "OP26_Analytics_Presentation.pptx")
    prs.save(output_path_app)
    print(f"Presentation copy saved successfully to {output_path_app}")

    # Save to deck/
    deck_dir = os.path.join(project_dir, "deck")
    os.makedirs(deck_dir, exist_ok=True)
    output_path_deck = os.path.join(deck_dir, "OP26_Analytics_Presentation.pptx")
    prs.save(output_path_deck)
    print(f"Presentation copy saved successfully to {output_path_deck}")
if __name__ == '__main__':
    create_presentation()
